from pathlib import Path
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "slides" / "pycycle_ai_prompting_course.pptx"
ASSET_DIR = ROOT / "docs" / "slides" / "assets"

FONT = "Microsoft JhengHei"
TITLE = RGBColor(74, 48, 38)
TEXT = RGBColor(53, 45, 38)
MUTED = RGBColor(112, 98, 84)
TEAL = RGBColor(40, 126, 116)
ORANGE = RGBColor(184, 93, 54)
GOLD = RGBColor(194, 142, 52)
BLUE = RGBColor(67, 102, 151)
GREEN = RGBColor(91, 138, 91)
RED = RGBColor(166, 82, 70)
LAV = RGBColor(121, 92, 145)
LIGHT = RGBColor(249, 242, 230)
LINE = RGBColor(218, 202, 180)
DARK = RGBColor(48, 39, 32)
BG = RGBColor(246, 238, 224)
CARD = RGBColor(255, 249, 239)
PATTERN = RGBColor(231, 214, 190)


def rgb(hex_value):
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_font(run, size=18, bold=False, color=TEXT):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        set_font(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.16, 12.2, 0.68, title, size=29, bold=True, color=TITLE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.88), Inches(1.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.55, 0.99, 12.2, 0.32, subtitle, size=14, color=MUTED)


def add_footer(slide, n):
    add_textbox(slide, 0.55, 7.10, 6.8, 0.28, "pyCycle + AI Prompting + uv + PySide6 教學規劃", size=10, color=MUTED)
    add_textbox(slide, 12.28, 7.10, 0.55, 0.28, str(n), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def bullet_box(slide, x, y, w, h, title, bullets, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    add_textbox(slide, x + 0.2, y + 0.15, w - 0.4, 0.42, title, size=18, bold=True, color=color)
    tx = "\n".join(f"• {b}" for b in bullets)
    add_textbox(slide, x + 0.25, y + 0.68, w - 0.45, h - 0.82, tx, size=14.2, color=TEXT)
    return shape


def prompt_card(slide, x, y, w, h, title, prompt, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = accent
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.36, title, size=15, bold=True, color=accent)
    add_textbox(slide, x + 0.18, y + 0.55, w - 0.36, h - 0.65, prompt, size=12.5, color=DARK)
    return shape


def pill(slide, x, y, w, h, text, fill, color=RGBColor(255, 255, 255), size=11):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    add_textbox(slide, x, y + 0.05, w, h - 0.02, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)
    return s


def arrow(slide, x1, y1, x2, y2, color=LINE, width=2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def safe_font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/msjhbd.ttc" if bold else "C:/Windows/Fonts/msjh.ttc",
        "C:/Windows/Fonts/mingliu.ttc",
        "C:/Windows/Fonts/arial.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


def pil_text(draw, xy, text, size=34, fill=(53, 45, 38), bold=False, anchor=None):
    draw.text(xy, text, font=safe_font(size, bold), fill=fill, anchor=anchor)


def make_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    palette = {
        "bg": (246, 238, 224),
        "card": (255, 249, 239),
        "ink": (53, 45, 38),
        "muted": (112, 98, 84),
        "line": (218, 202, 180),
        "orange": (184, 93, 54),
        "teal": (40, 126, 116),
        "blue": (67, 102, 151),
        "green": (91, 138, 91),
        "lav": (121, 92, 145),
        "red": (166, 82, 70),
        "gold": (194, 142, 52),
    }

    def canvas(name, title):
        img = Image.new("RGB", (1600, 900), palette["bg"])
        d = ImageDraw.Draw(img)
        for x in range(-200, 1700, 90):
            d.line((x, 0, x + 700, 900), fill=(236, 222, 201), width=2)
        for x in range(80, 1560, 160):
            for y in range(80, 860, 160):
                d.ellipse((x, y, x + 8, y + 8), fill=(225, 205, 178))
        d.rounded_rectangle((58, 54, 1542, 846), radius=32, fill=palette["card"], outline=palette["line"], width=4)
        pil_text(d, (95, 86), title, 46, palette["ink"], True)
        return img, d

    img, d = canvas("turbofan_cutaway.png", "高旁通比渦扇概念圖")
    d.rounded_rectangle((150, 355, 1420, 560), radius=95, outline=palette["ink"], width=8)
    d.rounded_rectangle((235, 392, 1300, 520), radius=60, outline=palette["line"], width=5)
    sections = [
        ("進氣", 180, palette["blue"]), ("風扇", 345, palette["teal"]), ("分流", 505, palette["gold"]),
        ("壓縮機", 670, palette["lav"]), ("燃燒", 875, palette["orange"]), ("渦輪", 1065, palette["red"]),
        ("噴嘴", 1260, palette["green"]),
    ]
    for label, cx, col in sections:
        d.ellipse((cx - 50, 365, cx + 50, 550), fill=col, outline=palette["ink"], width=3)
        pil_text(d, (cx, 600), label, 31, palette["ink"], True, anchor="mm")
    d.line((430, 330, 1330, 225), fill=palette["teal"], width=10)
    d.line((430, 585, 1330, 690), fill=palette["red"], width=10)
    pil_text(d, (705, 205), "旁通流：BPR、Fg、TSFC", 32, palette["teal"], True)
    pil_text(d, (705, 720), "核心流：壓縮、燃燒、渦輪功平衡", 32, palette["red"], True)
    img.save(ASSET_DIR / "turbofan_cutaway.png")

    img, d = canvas("ai_workflow.png", "AI 協作工作流")
    steps = [("規則", "AGENTS.md"), ("環境", "uv .venv"), ("資料", "Reference_sources"), ("模型", "pyCycle"), ("介面", "PySide6"), ("報告", "中 / EN")]
    xs = [180, 425, 670, 915, 1160, 1405]
    colors = [palette["blue"], palette["teal"], palette["orange"], palette["lav"], palette["green"], palette["gold"]]
    for i, ((a, b), x, col) in enumerate(zip(steps, xs, colors)):
        d.rounded_rectangle((x - 85, 345, x + 85, 505), radius=28, fill=col, outline=palette["ink"], width=3)
        pil_text(d, (x, 382), a, 34, (255, 255, 255), True, anchor="mm")
        pil_text(d, (x, 438), b, 24, (255, 255, 255), False, anchor="mm")
        if i < len(xs) - 1:
            d.line((x + 95, 425, xs[i + 1] - 100, 425), fill=palette["ink"], width=5)
            d.polygon([(xs[i + 1] - 100, 425), (xs[i + 1] - 122, 410), (xs[i + 1] - 122, 440)], fill=palette["ink"])
    pil_text(d, (130, 650), "關鍵：每一步都要求 AI 回報命令、檔案、驗證證據", 35, palette["ink"], True)
    img.save(ASSET_DIR / "ai_workflow.png")

    img, d = canvas("uv_environment.png", "uv 管理 .venv")
    d.rounded_rectangle((160, 250, 660, 650), radius=28, fill=(35, 35, 35), outline=palette["ink"], width=4)
    cmd_lines = ["uv init --bare", "uv add om-pycycle pyside6", "uv run python app.py", ".venv / pyproject.toml"]
    for i, line in enumerate(cmd_lines):
        pil_text(d, (205, 315 + i * 72), f"> {line}", 30, (245, 242, 230), False)
    d.rounded_rectangle((850, 255, 1320, 650), radius=36, fill=(255, 255, 255), outline=palette["line"], width=5)
    pil_text(d, (1085, 320), "可重現", 44, palette["teal"], True, anchor="mm")
    pil_text(d, (1085, 405), "不污染全域 Python", 34, palette["ink"], True, anchor="mm")
    pil_text(d, (1085, 485), "學生電腦可照做", 34, palette["ink"], True, anchor="mm")
    pil_text(d, (1085, 565), "錯誤可追蹤", 34, palette["ink"], True, anchor="mm")
    img.save(ASSET_DIR / "uv_environment.png")

    img, d = canvas("ui_mockup_large.png", "PySide6 教學 UI 雛形")
    d.rounded_rectangle((140, 210, 1470, 720), radius=30, fill=(255, 255, 255), outline=palette["line"], width=4)
    d.rectangle((140, 210, 1470, 285), fill=(238, 225, 207))
    pil_text(d, (175, 238), "Engine Deck UI", 34, palette["ink"], True)
    labels = ["Mach", "Altitude", "Fn target", "BPR", "T4"]
    for i, lab in enumerate(labels):
        y = 345 + i * 60
        pil_text(d, (190, y), lab, 28, palette["muted"], True)
        d.rounded_rectangle((380, y - 15, 590, y + 28), radius=10, fill=palette["card"], outline=palette["line"], width=2)
    d.rounded_rectangle((185, 655, 340, 705), radius=16, fill=palette["teal"])
    pil_text(d, (262, 681), "Run", 28, (255, 255, 255), True, anchor="mm")
    d.rounded_rectangle((370, 655, 555, 705), radius=16, fill=palette["orange"])
    pil_text(d, (462, 681), "Report", 28, (255, 255, 255), True, anchor="mm")
    rows = ["Fn 5900 lbf", "TSFC 0.56", "OPR 28.5", "FAR 0.024", "BPR 5.10"]
    for i, row in enumerate(rows):
        y = 355 + i * 58
        d.rounded_rectangle((760, y - 20, 1260, y + 25), radius=12, fill=(249, 242, 230) if i % 2 else palette["card"], outline=palette["line"], width=1)
        pil_text(d, (790, y - 5), row, 28, palette["ink"], True)
    img.save(ASSET_DIR / "ui_mockup_large.png")

    img, d = canvas("report_bilingual.png", "中英文工程報告")
    d.rounded_rectangle((160, 225, 725, 705), radius=24, fill=(255, 255, 255), outline=palette["line"], width=4)
    d.rounded_rectangle((875, 225, 1440, 705), radius=24, fill=(255, 255, 255), outline=palette["line"], width=4)
    pil_text(d, (205, 275), "正體中文報告", 38, palette["teal"], True)
    pil_text(d, (920, 275), "English Report", 38, palette["blue"], True)
    zh = ["摘要", "輸入條件", "主要結果", "來源比較", "模型限制"]
    en = ["Executive Summary", "Inputs", "Key Results", "Source Comparison", "Model Limits"]
    for i, line in enumerate(zh):
        pil_text(d, (220, 350 + i * 58), f"• {line}", 31, palette["ink"])
    for i, line in enumerate(en):
        pil_text(d, (940, 350 + i * 58), f"• {line}", 31, palette["ink"])
    img.save(ASSET_DIR / "report_bilingual.png")

    img, d = canvas("source_board.png", "Reference_sources 資料板")
    cards = [
        ("官方 URL", "CFM / Safran", palette["blue"]),
        ("PDF", "規格與論文", palette["orange"]),
        ("Manifest", "URL / 日期 / 欄位", palette["teal"]),
        ("驗證", "可對照 pyCycle", palette["green"]),
    ]
    for i, (a, b, col) in enumerate(cards):
        x = 170 + i * 345
        d.rounded_rectangle((x, 310, x + 260, 590), radius=28, fill=col, outline=palette["ink"], width=3)
        pil_text(d, (x + 130, 390), a, 36, (255, 255, 255), True, anchor="mm")
        pil_text(d, (x + 130, 470), b, 27, (255, 255, 255), False, anchor="mm")
    pil_text(d, (150, 690), "資料先存、來源先記、再進模型驗證", 42, palette["ink"], True)
    img.save(ASSET_DIR / "source_board.png")


def theme_slide(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    for i in range(7):
        x = 11.0 + i * 0.22
        y = 0.12 + i * 0.11
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.12), Inches(0.12))
        shp.fill.solid()
        shp.fill.fore_color.rgb = PATTERN
        shp.line.fill.background()


def add_picture(slide, name, x, y, w, h=None):
    path = ASSET_DIR / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))

    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        final_w = w
        final_h = w / img_ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * img_ratio
        final_x = x + (w - final_w) / 2
        final_y = y
    return slide.shapes.add_picture(str(path), Inches(final_x), Inches(final_y), width=Inches(final_w), height=Inches(final_h))


def two_col(slide, left_title, left_bullets, right_title, right_bullets, colors=(BLUE, ORANGE)):
    bullet_box(slide, 0.75, 1.55, 5.85, 4.95, left_title, left_bullets, color=colors[0])
    bullet_box(slide, 6.85, 1.55, 5.85, 4.95, right_title, right_bullets, color=colors[1])


def add_table(slide, x, y, w, h, headers, rows, col_widths=None, font_size=9.5):
    font_size = max(font_size + 2.2, 11.8)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE
        cell.text = head
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_font(r, size=font_size, bold=True, color=RGBColor(255, 255, 255))
    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i % 2 else LIGHT
            cell.text = value
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    set_font(r, size=font_size, color=TEXT)
    return table


def engine_schematic(slide, x, y, scale=1.0):
    # Simple teaching schematic, not a measured engine drawing.
    parts = [
        ("進氣\nInlet", x, y + 0.55, 1.0, 0.75, BLUE),
        ("風扇\nFan", x + 1.05, y + 0.35, 0.95, 1.15, TEAL),
        ("分流\nSplitter", x + 2.05, y + 0.50, 0.95, 0.85, GOLD),
        ("壓縮機\nLPC/HPC", x + 3.05, y + 0.45, 1.35, 0.95, LAV),
        ("燃燒室\nBurner", x + 4.55, y + 0.50, 1.05, 0.85, ORANGE),
        ("渦輪\nHPT/LPT", x + 5.75, y + 0.45, 1.25, 0.95, RED),
        ("噴嘴\nNozzles", x + 7.15, y + 0.55, 1.25, 0.75, GREEN),
    ]
    for label, px, py, pw, ph, color in parts:
        pill(slide, px, py, pw * scale, ph * scale, label, color, size=9)
    arrow(slide, x + 0.95, y + 0.92, x + 7.05, y + 0.92, color=MUTED, width=2)
    arrow(slide, x + 2.55, y + 0.50, x + 7.45, y + 0.15, color=TEAL, width=2)
    add_textbox(slide, x + 2.75, y + 0.02, 4.7, 0.25, "旁通流 bypass stream：提供大部分推力，影響 BPR / Fg / TSFC", size=8.5, color=TEAL)
    add_textbox(slide, x + 2.75, y + 1.45, 4.7, 0.25, "核心流 core stream：壓縮、燃燒、渦輪功平衡", size=8.5, color=RED)


def ui_mockup(slide, x, y, w, h):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 251, 252)
    bg.line.color.rgb = LINE
    add_textbox(slide, x + 0.25, y + 0.15, w - 0.5, 0.35, "PySide6 Engine Deck UI 雛形", size=14, bold=True, color=TITLE)
    labels = ["Mach", "Altitude ft", "Fn target lbf", "BPR", "T4 degR"]
    vals = ["0.80", "35000", "5900", "5.10", "2857"]
    for idx, (lab, val) in enumerate(zip(labels, vals)):
        yy = y + 0.7 + idx * 0.45
        add_textbox(slide, x + 0.3, yy, 1.5, 0.22, lab, size=8.5, color=MUTED)
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1.55), Inches(yy - 0.02), Inches(1.1), Inches(0.28))
        r.fill.solid()
        r.fill.fore_color.rgb = RGBColor(255, 255, 255)
        r.line.color.rgb = LINE
        add_textbox(slide, x + 1.65, yy, 0.9, 0.2, val, size=8.5, color=TEXT)
    pill(slide, x + 0.3, y + 3.08, 1.1, 0.34, "Run", TEAL, size=9)
    pill(slide, x + 1.55, y + 3.08, 1.1, 0.34, "Report", ORANGE, size=9)
    rows = [
        ("Fn", "5900 lbf"),
        ("TSFC", "0.56 lbm/hr/lbf"),
        ("OPR", "28.5"),
        ("FAR", "0.024"),
        ("BPR", "5.10"),
    ]
    add_table(slide, x + 3.0, y + 0.65, w - 3.35, 2.8, ["輸出", "值"], rows, col_widths=[1.2, 1.55], font_size=8.3)


def make_deck():
    make_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []

    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    add_textbox(s, 0.75, 0.55, 8.8, 1.1, "用 AI 提示詞、uv、pyCycle 與 PySide6 建立渦扇引擎教學 UI", size=32, bold=True, color=TITLE)
    add_textbox(s, 0.78, 1.78, 8.8, 0.65, "以 CFM56-7B 等級高旁通比渦扇為案例，從資料治理、模型驗證到中英文報告產出", size=17, color=MUTED)
    add_picture(s, "turbofan_cutaway.png", 0.72, 2.65, 7.7, 4.15)
    prompt_card(s, 9.35, 0.85, 3.1, 5.4, "本課核心任務", "學生要學會：\n1. 用清楚提示詞指揮 AI\n2. 用 uv 管理 .venv\n3. 不修改 upstream/pyCycle\n4. 建立 PySide6 UI\n5. 產出正體中文與英文工程報告", accent=ORANGE)
    slides.append(s)

    data = [
        ("1", "提示詞與規則", "讀 AGENTS.md、建立 AI 工作邊界"),
        ("2", "uv 環境", "建立 .venv、安裝 pyCycle / PySide6"),
        ("3", "參考資料", "CFM56-7B 來源存入 Reference_sources"),
        ("4", "pyCycle 驗證", "跑 high_bypass_turbofan 範例並讀輸出"),
        ("5", "PySide6 UI", "建立輸入面板、結果表、報告按鈕"),
        ("6", "報告與延伸", "正體中文/英文報告、參數掃描、GasTurb 對照"),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "課程地圖", "六個階段，每階段都有可複製提示詞與驗證產出")
    for i, (num, title, desc) in enumerate(data):
        x = 0.8 + (i % 3) * 4.15
        y = 1.6 + (i // 3) * 2.25
        pill(s, x, y, 0.55, 0.55, num, [BLUE, TEAL, ORANGE, LAV, GREEN, RED][i], size=15)
        bullet_box(s, x + 0.7, y - 0.08, 3.2, 1.45, title, [desc], color=[BLUE, TEAL, ORANGE, LAV, GREEN, RED][i])
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "AI 提示詞設計原則", "先規則、再目標、再限制、再驗證")
    two_col(
        s,
        "好提示詞包含",
        ["角色：你是嚴謹的教學助教", "工作目標：建立 uv + pyCycle + UI", "限制：不得修改 upstream/pyCycle", "輸出格式：步驟、命令、驗證結果", "資料規則：先存 Reference_sources"],
        "避免這些寫法",
        ["只說「幫我做 UI」", "沒有指定資料來源與引用方式", "沒有要求 AI 先讀 repo 規則", "沒有要求驗證命令", "把真實引擎數值直接硬塞模型"],
        colors=(TEAL, RED),
    )
    slides.append(s)

    prompts = [
        ("啟動專案提示詞", "請先閱讀 D:\\45_pyCycle_edu\\AGENTS.md 與 docs/dev-log/README.md。\n任務：規劃 pyCycle 教學專案。\n限制：不要修改 upstream/pyCycle；任何資料來源先記錄到 Reference_sources。\n輸出：請列出你會建立的檔案、驗證命令、風險。"),
        ("要求 AI 先查規則", "在分析或修改之前，請從 repo 根目錄執行：npx gitnexus analyze --embeddings。\n完成後摘要 GitNexus 結果，並說明哪些檔案可以改、哪些只能讀。"),
        ("要求 AI 回報證據", "每完成一階段，請列出：\n1. 實際執行命令\n2. 產生或修改的檔案\n3. 驗證結果\n4. 下一步風險"),
    ]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 0：讓 AI 先遵守專案規則", "這張可直接給學生當第一個練習")
    for i, (t, p) in enumerate(prompts):
        prompt_card(s, 0.75 + i * 4.15, 1.45, 3.85, 4.9, t, p, accent=[TEAL, BLUE, ORANGE][i])
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 1：用 uv 管理 .venv", "把環境管理變成可重現的提示詞任務")
    prompt_card(s, 0.7, 1.35, 5.75, 5.2, "學生提示詞", "請在 D:\\45_pyCycle_edu 建立 uv 管理的 Python 環境。\n需求：\n- 建立 .venv\n- 安裝 om-pycycle、openmdao、pyside6、matplotlib、pandas\n- 產生 pyproject.toml\n- 用 uv run 驗證 import pycycle 與 import PySide6\n請先列命令，再執行，最後摘要結果。", accent=TEAL)
    add_picture(s, "uv_environment.png", 6.8, 1.45, 5.65, 3.2)
    bullet_box(s, 6.9, 4.95, 5.45, 1.25, "驗證重點", ["不能用全域 Python 冒充 .venv", "若套件衝突，請 AI 解釋版本原因"], color=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 2：參考資料治理", "CFM56-7B 的資料先存，再拿來驗證")
    add_table(
        s, 0.7, 1.35, 12.0, 2.45,
        ["資料", "用途", "存放方式"],
        [
            ["CFM / Safran 官方頁", "發動機系列背景、推力級距", "URL + 擷取欄位 + 下載狀態"],
            ["CFM56-7B specs PDF", "BPR、推力、尺寸等交叉檢查", "PDF 存 Reference_sources"],
            ["pyCycle README / paper", "模型能力與範例來源", "README / PDF 存 Reference_sources"],
            ["pyCycle high_bypass_turbofan.py", "教學模型骨架", "只讀 upstream，不修改"],
        ],
        col_widths=[3.2, 4.7, 4.1],
        font_size=9.2,
    )
    add_picture(s, "source_board.png", 0.85, 4.0, 4.65, 2.62)
    prompt_card(s, 5.75, 4.15, 6.75, 1.85, "資料蒐集提示詞", "請以 CFM56-7B 為主題搜尋官方或可信來源。所有資料先寫入 D:\\45_pyCycle_edu\\Reference_sources\\source_manifest.md。\n每筆資料需包含：URL、存取日期、下載狀態、可驗證欄位、是否能直接對照 pyCycle。", accent=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "為什麼第一版選 CFM56-7B", "主流、資料多、和 pyCycle 高旁通比範例相容")
    add_picture(s, "turbofan_cutaway.png", 0.75, 1.55, 6.0, 3.38)
    bullet_box(s, 7.05, 1.55, 5.25, 1.45, "教學優點", ["737NG 常見，學生容易理解", "公開資料多，可練習交叉驗證"], color=TEAL)
    bullet_box(s, 7.05, 3.2, 5.25, 1.45, "模型定位", ["pyCycle 範例是教學骨架", "先做趨勢與量級驗證"], color=BLUE)
    bullet_box(s, 7.05, 4.85, 5.25, 1.45, "第一版限制", ["不做完整 GasTurb 複刻", "不宣稱精準重現 CFM56-7B"], color=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 3：pyCycle 可以得到什麼資訊", "從模型輸出轉成學生看得懂的工程報告")
    rows = [
        ["整體性能", "Fn, Fg, ram drag, OPR, TSFC, BPR", "推力、油耗、旁通比、總壓比"],
        ["飛行條件", "Mach, altitude, dTs, inlet W", "不同高度/速度下的工作點"],
        ["流站資料", "Pt, Tt, ht, entropy, Ps, W, MN, V, area", "每一站總壓/總溫/質流/面積"],
        ["壓縮機", "Wc, PR, eta, Nc, map line, surge margin", "壓縮效率與喘振裕度"],
        ["燃燒室", "TtOut, Wfuel, FAR, pressure loss", "燃油流量與燃氣溫度"],
        ["渦輪/軸", "PR, eff, power, torque balance", "高低壓軸功率平衡"],
        ["噴嘴", "PR, throat area, MN, V, gross thrust", "核心與旁通噴嘴推力"],
    ]
    add_table(s, 0.55, 1.35, 12.25, 5.25, ["類別", "pyCycle 變數", "可教學意義"], rows, col_widths=[1.6, 5.0, 5.65], font_size=8.7)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "pyCycle 範例切入點", "high_bypass_turbofan.py 的課堂閱讀路線")
    add_table(
        s,
        0.75,
        1.35,
        11.85,
        3.65,
        ["區塊", "程式概念", "教學意義"],
        [
            ["Cycle", "HBTF(pyc.Cycle)", "建立高旁通比渦扇架構"],
            ["Elements", "Inlet, Compressor, Splitter, Combustor, Turbine, Nozzle", "每個元件代表一段熱力或流體模型"],
            ["Balances", "W, FAR, lpt_PR, hpt_PR, BPR, spool speed", "求解推力、T4、功率平衡與噴嘴面積"],
            ["Viewer", "summary + flow station + component tables", "把 OpenMDAO 結果轉成可讀報表"],
        ],
        col_widths=[1.7, 4.9, 5.25],
        font_size=9.3,
    )
    prompt_card(s, 0.9, 5.45, 11.5, 1.15, "閱讀提示詞", "請閱讀 upstream/pyCycle/example_cycles/high_bypass_turbofan.py，只做摘要，不修改檔案。請用表格說明每個子系統、balance 變數與 viewer 輸出。", accent=TEAL)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 4：跑範例與驗證", "先跑原始範例，再建立 wrapper 腳本")
    prompt_card(s, 0.75, 1.35, 5.9, 4.9, "執行提示詞", "請使用 uv run 執行 pyCycle high_bypass_turbofan 範例。\n限制：不要修改 upstream/pyCycle。\n請把輸出檔複製到 wrapper 的 results/ 目錄，並摘要 DESIGN 與 OD_part_pwr 的 Fn、TSFC、OPR、BPR、FAR。\n若失敗，請先解釋錯誤來源再修正環境。", accent=BLUE)
    bullet_box(s, 7.05, 1.35, 5.4, 2.05, "驗證命令", ["uv run python upstream/pyCycle/example_cycles/high_bypass_turbofan.py", "檢查 hbtf_view.out 是否產生", "比對 Reference_sources 的 BPR/推力量級"], color=TEAL)
    bullet_box(s, 7.05, 3.75, 5.4, 2.15, "學生要回答", ["設計點與非設計點差在哪？", "為什麼 off-design 要平衡噴嘴面積？", "TSFC 的單位與意義是什麼？", "哪些結果不能直接宣稱等於 CFM56-7B？"], color=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "階段 5：PySide6 UI 功能規劃", "先做工程工具，不做華麗展示頁")
    add_picture(s, "ui_mockup_large.png", 0.75, 1.45, 6.0, 3.38)
    bullet_box(s, 7.2, 1.25, 5.25, 1.58, "第一版輸入", ["Mach、altitude、dTs", "Fn target、BPR、T4", "設計點/非設計點選擇"], color=BLUE)
    bullet_box(s, 7.2, 3.0, 5.25, 1.58, "第一版輸出", ["性能摘要表", "流站與元件表", "收斂狀態與警告"], color=TEAL)
    bullet_box(s, 7.2, 4.75, 5.25, 1.58, "第一版報告", ["正體中文報告", "English report", "引用來源與參數表"], color=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "UI 建置提示詞", "要求 AI 先設計資料流，再寫 PySide6")
    prompt_card(s, 0.7, 1.3, 5.95, 5.25, "介面設計提示詞", "請設計一個 PySide6 UI，用於操作 pyCycle 高旁通比渦扇教學模型。\n限制：不要修改 upstream/pyCycle；模型呼叫放在 wrapper 模組。\n畫面需包含：輸入表單、Run 按鈕、結果表、收斂訊息、匯出報告按鈕。\n請先提出檔案結構與資料流，不要直接寫程式。", accent=TEAL)
    prompt_card(s, 6.95, 1.3, 5.75, 5.25, "實作提示詞", "根據前一步設計，請建立最小可用 PySide6 應用。\n需求：\n- 用 QMainWindow + QWidget\n- 輸入值有驗證\n- 長時間運算不可卡住 UI\n- 結果用 QTableWidget 顯示\n- 報告輸出為 Markdown\n請附上 uv run 啟動命令與測試方式。", accent=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "GasTurb 功能趨近路線", "不是複刻介面，而是逐步補齊教學所需能力")
    rows = [
        ["1", "單一工作點", "輸入 Mach/altitude/Fn，輸出性能表", "pyCycle example + UI"],
        ["2", "參數掃描", "改變 BPR/T4/高度，看 TSFC 與 Fn 趨勢", "pandas + matplotlib"],
        ["3", "流程圖視覺化", "顯示 inlet/fan/core/bypass/nozzle", "PySide6 scene 或簡圖"],
        ["4", "資料驗證", "把來源數值與模型結果並列", "Reference_sources manifest"],
        ["5", "報告產生", "中文/英文工程摘要與引用", "Markdown / PDF"],
    ]
    add_table(s, 0.75, 1.4, 11.85, 4.65, ["階段", "功能", "學生看到的價值", "技術"], rows, col_widths=[0.8, 2.4, 5.1, 3.55], font_size=9)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "報告產生：正體中文與英文", "同一份模型結果，輸出兩種工程語言")
    add_picture(s, "report_bilingual.png", 0.85, 1.42, 5.6, 3.15)
    prompt_card(s, 6.7, 1.35, 5.75, 2.25, "中文報告提示詞", "請根據 pyCycle 結果產生正體中文工程報告。\n章節：摘要、輸入條件、主要結果、與 CFM56-7B 公開資料比較、限制與下一步。\n要求：所有數值附單位；引用 Reference_sources/source_manifest.md；不要誇大模型精度。", accent=TEAL)
    prompt_card(s, 6.7, 3.95, 5.75, 2.2, "English report prompt", "Generate an engineering report in English from the pyCycle output.\nSections: Executive Summary, Inputs, Key Results, Comparison Against Public CFM56-7B Data, Model Limitations, Next Steps.\nUse units for every value and cite Reference_sources/source_manifest.md.", accent=BLUE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "學生作業設計", "每一週都有可驗證產出")
    rows = [
        ["第 1 週", "提示詞與 repo 規則", "能讓 AI 正確遵守不改 upstream"],
        ["第 2 週", "uv + pyCycle 範例", "能重現 high_bypass_turbofan 輸出"],
        ["第 3 週", "CFM56-7B 資料表", "能說明每個來源能驗證什麼"],
        ["第 4 週", "PySide6 最小 UI", "能輸入條件並顯示結果"],
        ["第 5 週", "中英文報告", "能引用來源並說明限制"],
        ["第 6 週", "參數掃描展示", "能用圖表說明 BPR/T4/高度趨勢"],
    ]
    add_table(s, 0.8, 1.25, 11.7, 5.15, ["週次", "主題", "通過標準"], rows, col_widths=[1.2, 3.6, 6.9], font_size=10)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "評分規準", "評估 AI 使用能力，不只看程式跑不跑")
    two_col(
        s,
        "技術面",
        ["uv 環境可重現", "pyCycle 結果可解釋", "UI 不阻塞且有輸入驗證", "報告數值有單位", "資料來源可追溯"],
        "AI 協作面",
        ["提示詞清楚且可複製", "AI 有先讀規則", "錯誤修正有原因分析", "沒有未驗證宣稱", "有記錄 dev-log"],
        colors=(BLUE, ORANGE),
    )
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "常見失敗與修正提示詞", "把錯誤當成提示詞改寫訓練")
    rows = [
        ["套件裝不起來", "請檢查 pyproject.toml 與 uv lock，說明版本衝突來源，提出最小修正。"],
        ["pyCycle 不收斂", "請不要任意改模型，先列出 balance 變數、初始猜測與可能物理原因。"],
        ["UI 卡住", "請把模型運算移到 worker thread，並讓 UI 顯示 running / failed / completed 狀態。"],
        ["資料來源混亂", "請重新整理 Reference_sources/source_manifest.md，只保留可追溯來源。"],
        ["報告太像廣告", "請改寫為工程報告語氣，列出限制，不要宣稱等同原廠資料。"],
    ]
    add_table(s, 0.75, 1.35, 11.8, 4.95, ["問題", "修正提示詞"], rows, col_widths=[2.3, 9.5], font_size=9.2)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "課堂 Demo 腳本", "老師可以照這張投影片帶一次完整流程")
    add_picture(s, "ai_workflow.png", 0.95, 1.45, 11.2, 2.35)
    steps = [
        ("01", "要求 AI 讀 AGENTS.md"),
        ("02", "用 uv 建環境"),
        ("03", "整理 CFM56-7B 來源"),
        ("04", "跑 pyCycle 範例"),
        ("05", "解讀 Fn / TSFC / OPR / BPR"),
        ("06", "開 PySide6 UI"),
        ("07", "輸出中英文報告"),
    ]
    for i, (num, text) in enumerate(steps):
        x = 0.7 + i * 1.78
        pill(s, x, 4.05, 0.65, 0.55, num, [BLUE, TEAL, ORANGE, LAV, GREEN, RED, GOLD][i], size=13)
        add_textbox(s, x - 0.2, 4.78, 1.45, 0.9, text, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(s, x + 0.72, 4.32, x + 1.45, 4.32, color=MUTED, width=1.5)
    prompt_card(s, 1.05, 6.05, 11.0, 0.75, "總結提示詞", "請回顧今天所有步驟，產生一份學習紀錄：包含使用過的提示詞、成功命令、失敗與修正、目前模型能輸出的 pyCycle 資訊、下一週要做的 UI 功能。", accent=LAV)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "檔案結構建議", "後續實作時維持 wrapper / upstream 分離")
    rows = [
        ["Reference_sources/", "來源、URL、PDF、擷取欄位、下載狀態"],
        ["src/pycycle_edu/", "wrapper 模型、資料轉換、報告產生"],
        ["src/pycycle_edu_ui/", "PySide6 視窗、worker、結果表"],
        ["results/", "pyCycle 執行輸出、CSV、圖表"],
        ["reports/", "正體中文與英文報告"],
        ["upstream/pyCycle/", "唯讀參考，不修改"],
        ["docs/dev-log/", "每次 meaningful task 的紀錄"],
    ]
    add_table(s, 0.85, 1.35, 11.45, 4.8, ["路徑", "用途"], rows, col_widths=[3.0, 8.45], font_size=10)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "教師提醒", "讓學生知道 AI 很會做事，但要被清楚管理")
    bullet_box(s, 0.85, 1.45, 3.6, 4.6, "工程倫理", ["公開資料不等於可隨意宣稱精準", "模型假設要寫清楚", "報告需要引用來源"], color=TEAL)
    bullet_box(s, 4.85, 1.45, 3.6, 4.6, "軟體習慣", ["每次先讀規則", "用 uv 重現環境", "錯誤要有紀錄", "不要改 upstream"], color=BLUE)
    bullet_box(s, 8.85, 1.45, 3.6, 4.6, "AI 使用", ["提示詞要具體", "要求 AI 給驗證證據", "把大型任務拆階段", "保留失敗案例"], color=ORANGE)
    slides.append(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "下一步：從 PPT 到可執行教材", "本投影片是規劃與教學腳本，下一階段才進入程式實作")
    prompt_card(s, 0.85, 1.4, 11.65, 2.0, "下一階段給 AI 的提示詞", "請根據 docs/slides/pycycle_ai_prompting_course.pptx，建立第一版可執行教材。\n範圍：uv 專案設定、pyCycle wrapper、PySide6 最小 UI、Markdown 報告輸出。\n限制：不要修改 upstream/pyCycle；修改任何既有 symbol 前先做 GitNexus impact analysis；每階段更新 docs/dev-log。", accent=TEAL)
    bullet_box(s, 0.95, 4.05, 3.55, 1.65, "第一個可執行里程碑", ["uv sync", "跑 high_bypass_turbofan", "輸出 results/hbtf_summary.csv"], color=BLUE)
    bullet_box(s, 4.95, 4.05, 3.55, 1.65, "第二個里程碑", ["PySide6 表單", "Run 按鈕", "結果表"], color=ORANGE)
    bullet_box(s, 8.95, 4.05, 3.55, 1.65, "第三個里程碑", ["中文報告", "英文報告", "來源引用"], color=GREEN)
    slides.append(s)

    max_slides = os.environ.get("PYCE_PPT_MAX_SLIDES")
    if max_slides:
        keep = int(max_slides)
        while len(prs.slides) > keep:
            r_id = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[-1]
        slides = slides[:keep]

    for idx, slide in enumerate(slides, start=1):
        theme_slide(slide)
        add_footer(slide, idx)

    out = Path(os.environ.get("PYCE_PPT_OUT", OUT))
    prs.save(out)
    print(f"Generated {out}")


if __name__ == "__main__":
    make_deck()
